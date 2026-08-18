"""Pluggable text extraction (Strategy + Registry).

Each file type gets its own :class:`ContentExtractor`. A registry maps an
extension to the strategy that handles it, so adding EPUB support later means
writing one class and registering it -- no existing code changes. That is the
Open/Closed Principle applied to a dispatch table instead of an if/elif chain.

Every extractor is defensive: a corrupt PDF in the middle of a 10k-file scan
must degrade to "no text for this file", never abort the job.
"""

from __future__ import annotations

import csv
import io
import logging
import zipfile
from pathlib import Path
from typing import Protocol

from app.core.textutils import collapse_whitespace, truncate

logger = logging.getLogger(__name__)

# Read-cap for pathological inputs (a 500k-row CSV, a 4k-page PDF).
_MAX_PDF_PAGES = 80
_MAX_SHEET_ROWS = 500
_MAX_SLIDES = 120


class ContentExtractor(Protocol):
    """Strategy interface. ``kind`` is stored as searchable metadata.

    Structural, not nominal: an implementation just needs the two attributes
    and the method, with no base class to inherit.
    """

    kind: str
    extensions: frozenset[str]

    def extract(self, path: Path, max_chars: int) -> str: ...


class PlainTextExtractor:
    """Handles anything that is already UTF-8-ish text, including source code.

    Decoding strategy: try UTF-8 first (the overwhelmingly common case, and a
    strict codec that fails fast on non-UTF-8 bytes), then fall back to
    ``charset_normalizer`` which scores candidate encodings by how plausible
    the decoded text looks. Final fallback is lossy latin-1, which never
    raises.
    """

    kind = "text"
    extensions = frozenset(
        {
            ".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv",
            ".json", ".jsonl", ".yaml", ".yml", ".toml", ".ini", ".cfg",
            ".conf", ".env", ".xml", ".html", ".htm", ".css", ".scss",
            ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".kt", ".go",
            ".rs", ".rb", ".php", ".c", ".h", ".cpp", ".hpp", ".cs",
            ".sql", ".sh", ".bash", ".ps1", ".bat", ".r", ".m", ".swift",
            ".tex", ".gradle", ".dockerfile", ".gitignore", ".tf",
        }
    )

    def extract(self, path: Path, max_chars: int) -> str:
        # Bounded read rather than read_bytes()[:n], which would pull the whole
        # file into memory first. 4 bytes/char is the UTF-8 worst case, so this
        # can never truncate below the requested character budget.
        with open(path, "rb") as handle:
            raw = handle.read(max_chars * 4)
        if not raw:
            return ""
        text = self._decode(raw)
        if path.suffix.lower() in {".csv", ".tsv"}:
            text = self._flatten_delimited(text, path.suffix.lower())
        return truncate(collapse_whitespace(text), max_chars)

    @staticmethod
    def _decode(raw: bytes) -> str:
        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            pass
        try:
            from charset_normalizer import from_bytes

            best = from_bytes(raw).best()
            if best is not None:
                return str(best)
        except Exception:  # noqa: BLE001 - detection is strictly best-effort
            logger.debug("charset detection failed", exc_info=True)
        return raw.decode("latin-1", errors="replace")

    @staticmethod
    def _flatten_delimited(text: str, suffix: str) -> str:
        """Turn tabular rows into ``col: value`` prose.

        Embedding models are trained on sentences, not grids. Naming the
        column beside each value gives the vector far more to work with than a
        bare comma-separated line.
        """
        delimiter = "\t" if suffix == ".tsv" else ","
        try:
            reader = csv.reader(io.StringIO(text), delimiter=delimiter)
            rows = [row for _, row in zip(range(_MAX_SHEET_ROWS + 1), reader)]
        except csv.Error:
            return text
        if not rows:
            return ""
        header = [cell.strip() for cell in rows[0]]
        lines = [" | ".join(header)]
        for row in rows[1:]:
            pairs = [
                f"{header[i]}: {cell.strip()}"
                for i, cell in enumerate(row)
                if i < len(header) and cell.strip()
            ]
            if pairs:
                lines.append("; ".join(pairs))
        return "\n".join(lines)


class PdfExtractor:
    kind = "pdf"
    extensions = frozenset({".pdf"})

    def extract(self, path: Path, max_chars: int) -> str:
        from pypdf import PdfReader

        reader = PdfReader(str(path), strict=False)
        if reader.is_encrypted:
            try:
                reader.decrypt("")  # many PDFs are "encrypted" with no password
            except Exception:  # noqa: BLE001
                return ""

        parts: list[str] = []
        budget = max_chars
        for page_number, page in enumerate(reader.pages[:_MAX_PDF_PAGES], start=1):
            if budget <= 0:
                break
            try:
                page_text = page.extract_text() or ""
            except Exception:  # noqa: BLE001 - one bad page shouldn't kill the doc
                logger.debug("page %s of %s failed", page_number, path, exc_info=True)
                continue
            page_text = collapse_whitespace(page_text)
            if page_text:
                parts.append(f"[page {page_number}] {page_text}")
                budget -= len(page_text)
        return truncate("\n\n".join(parts), max_chars)


class DocxExtractor:
    kind = "docx"
    extensions = frozenset({".docx"})

    def extract(self, path: Path, max_chars: int) -> str:
        import docx

        document = docx.Document(str(path))
        parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return truncate(collapse_whitespace("\n".join(parts)), max_chars)


class XlsxExtractor:
    kind = "spreadsheet"
    extensions = frozenset({".xlsx", ".xlsm"})

    def extract(self, path: Path, max_chars: int) -> str:
        from openpyxl import load_workbook

        # read_only streams rows instead of building the full object graph;
        # data_only returns cached formula *results* rather than "=SUM(A1:A9)".
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        try:
            parts: list[str] = []
            for sheet in workbook.worksheets:
                parts.append(f"[sheet: {sheet.title}]")
                for index, row in enumerate(sheet.iter_rows(values_only=True)):
                    if index >= _MAX_SHEET_ROWS:
                        break
                    cells = [str(v).strip() for v in row if v is not None]
                    if cells:
                        parts.append(" | ".join(cells))
            return truncate(collapse_whitespace("\n".join(parts)), max_chars)
        finally:
            workbook.close()


class PptxExtractor:
    kind = "presentation"
    extensions = frozenset({".pptx"})

    def extract(self, path: Path, max_chars: int) -> str:
        from pptx import Presentation

        presentation = Presentation(str(path))
        parts: list[str] = []
        for number, slide in enumerate(presentation.slides, start=1):
            if number > _MAX_SLIDES:
                break
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            if texts:
                parts.append(f"[slide {number}] " + " ".join(texts))
        return truncate(collapse_whitespace("\n".join(parts)), max_chars)


class NullExtractor:
    """Terminal fallback for binaries: index metadata, never content."""

    kind = "binary"
    extensions = frozenset()

    def extract(self, path: Path, max_chars: int) -> str:  # noqa: ARG002
        return ""


class ExtractorRegistry:
    """Extension -> strategy dispatch table. Lookup is O(1)."""

    def __init__(self, extractors: list[ContentExtractor] | None = None) -> None:
        self._by_extension: dict[str, ContentExtractor] = {}
        self._fallback = NullExtractor()
        for extractor in extractors or self._defaults():
            self.register(extractor)

    @staticmethod
    def _defaults() -> list[ContentExtractor]:
        return [
            PlainTextExtractor(),
            PdfExtractor(),
            DocxExtractor(),
            XlsxExtractor(),
            PptxExtractor(),
        ]

    def register(self, extractor: ContentExtractor) -> None:
        for extension in extractor.extensions:
            self._by_extension[extension.lower()] = extractor

    def supports(self, extension: str) -> bool:
        return extension.lower() in self._by_extension

    def kind_for(self, extension: str) -> str:
        return self._by_extension.get(extension.lower(), self._fallback).kind

    def extract(self, path: Path, max_chars: int) -> tuple[str, str]:
        """Return ``(text, kind)``; ``text`` is empty when unsupported.

        Exceptions are swallowed by design -- a single unreadable file must
        not abort a long-running index job.
        """
        extractor = self._by_extension.get(path.suffix.lower(), self._fallback)
        try:
            return extractor.extract(path, max_chars), extractor.kind
        except (OSError, zipfile.BadZipFile, ValueError, KeyError, IndexError):
            logger.warning("extraction failed for %s", path, exc_info=True)
            return "", extractor.kind
        except Exception:  # noqa: BLE001 - third-party parsers raise anything
            logger.warning("unexpected extractor error for %s", path, exc_info=True)
            return "", extractor.kind
